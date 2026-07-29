"""Embedding API 服务 — FastAPI + BGE-M3，监听 0.0.0.0:5000"""
import os
os.environ.setdefault("HF_ENDPOINT", "https://hf-mirror.com")

import asyncio
import base64
import io
import re
import urllib.request
from typing import Optional

import pdfplumber
from fastapi import FastAPI
from fastapi.responses import JSONResponse
from pptx import Presentation
from docx import Document
from pydantic import BaseModel
from sentence_transformers import SentenceTransformer

app = FastAPI(title="Embedding API")

print("加载 BGE-M3 模型（冷启动约 10-30s）...")
model = SentenceTransformer("BAAI/bge-m3")
DIM = model.get_embedding_dimension()
print(f"模型就绪，向量维度: {DIM}")

EXTRACT_TEXT_TIMEOUT_SECONDS = 60
MAX_EXTRACTED_CHARS = 2000
MAX_EXTRACT_FILE_SIZE = 10 * 1024 * 1024

PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
XLS_MIME = "application/vnd.ms-excel"
PDF_MIME = "application/pdf"
TEXT_MIME_PREFIXES = (
    "text/plain",
    "text/markdown",
    "text/csv",
    "application/json",
    "application/xml",
    "text/xml",
)
DATA_URL_PATTERN = re.compile(r"^data:([^;,]+)?(;base64)?,(.*)$", re.DOTALL)


@app.get("/")
async def root():
    return {"status": "ok", "model": "BAAI/bge-m3", "dimension": DIM}


@app.get("/health")
async def health():
    return {"status": "healthy"}


@app.get("/dimension")
async def dimension():
    return {"dimension": DIM}


class EmbedRequest(BaseModel):
    text: str


class BatchEmbedRequest(BaseModel):
    texts: list[str]


@app.post("/embed")
async def embed(body: EmbedRequest):
    """接收单个文本，返回 1024 维向量（JSON 数组）"""
    emb = model.encode(body.text).tolist()
    return JSONResponse({"embedding": emb})


@app.post("/embed_batch")
async def embed_batch(body: BatchEmbedRequest):
    """批量编码，避免频繁调用开销"""
    embs = model.encode(body.texts).tolist()
    return JSONResponse({"embeddings": embs})


class ExtractRequest(BaseModel):
    url: str
    mimeType: str
    name: str
    size: int = 0
    page_from: int | None = None
    page_to: int | None = None
    max_chars: int | None = None


def _decode_data_url(url: str) -> bytes:
    match = DATA_URL_PATTERN.match(url)
    if not match:
        raise ValueError("invalid_data_url")
    _, is_base64, payload = match.groups()
    if is_base64:
        return base64.b64decode(payload)
    from urllib.parse import unquote

    return unquote(payload).encode("utf-8")


def _fetch_http_bytes(url: str) -> bytes:
    request = urllib.request.Request(url, headers={"User-Agent": "project-manager/extract"})
    with urllib.request.urlopen(request, timeout=EXTRACT_TEXT_TIMEOUT_SECONDS) as response:
        return response.read(MAX_EXTRACT_FILE_SIZE + 1)


def _read_bytes(url: str) -> bytes:
    if url.startswith("data:"):
        return _decode_data_url(url)
    if url.startswith("http://") or url.startswith("https://") or url.startswith("/"):
        return _fetch_http_bytes(url)
    raise ValueError("unsupported_url_scheme")


def _extract_text_from_bytes(
    raw: bytes,
    mime_type: str,
    page_from: int | None = None,
    page_to: int | None = None,
    name: str = "",
) -> str:
    normalized = (mime_type or "").split(";")[0].strip().lower()

    if normalized.startswith(TEXT_MIME_PREFIXES):
        return raw.decode("utf-8", errors="replace")

    if normalized == PDF_MIME:
        from extractors.pdf_ocr import extract_pdf_ocr
        return extract_pdf_ocr(raw, page_from=page_from, page_to=page_to)

    if normalized == PPTX_MIME:
        presentation = Presentation(io.BytesIO(raw))
        parts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        parts.append(text)
        return "\n".join(parts)

    if normalized == DOCX_MIME:
        doc = Document(io.BytesIO(raw))
        parts = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)
        for table in doc.tables:
            for row in table.rows:
                row_texts = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_texts.append(cell_text)
                if row_texts:
                    parts.append(" | ".join(row_texts))
        return "\n".join(parts)

    if normalized == XLSX_MIME:
        try:
            import openpyxl

            wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
            sheet_parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                row_lines = []
                for row in ws.iter_rows(values_only=True):
                    cell_values = []
                    for cell in row:
                        if cell is not None and cell != "":
                            cell_values.append(str(cell))
                    if cell_values:
                        row_lines.append("\t".join(cell_values))
                if row_lines:
                    sheet_parts.append(f"[{sheet_name}]\n" + "\n".join(row_lines))
            wb.close()
            return "\n\n".join(sheet_parts)
        except Exception as exc:
            raise ValueError(f"xlsx_parse_error:{exc}") from exc

    if normalized == XLS_MIME:
        raise ValueError("xls_not_supported")

    # 图片 OCR
    if normalized.startswith("image/"):
        from extractors.image import extract_image_text

        return extract_image_text(raw)

    # .doc (OLE 格式)
    if normalized == "application/msword" or name.lower().endswith(".doc"):
        from extractors.doc import extract_doc

        return extract_doc(raw)

    # WPS 文件
    if normalized in ("application/wpsoffice", "application/wps-office.docx") or name.lower().endswith(
        (".wps", ".wpt")
    ):
        from extractors.wps import extract_wps

        return extract_wps(raw)

    raise ValueError(f"unsupported_mime:{normalized}")


def _truncate(text: str, max_chars: int = MAX_EXTRACTED_CHARS) -> str:
    if len(text) <= max_chars:
        return text
    return f"{text[:max_chars].rstrip()}…"


def extract_attachment_text(
    url: str,
    mime_type: str,
    name: str,
    size: int,
    page_from: int | None = None,
    page_to: int | None = None,
    max_chars: int | None = None,
) -> dict:
    if size and size > MAX_EXTRACT_FILE_SIZE:
        return {"text": "", "source": "skipped_too_large"}

    try:
        raw = _read_bytes(url)
    except urllib.error.HTTPError as exc:
        print(f"[extract-text] storage not found for {name}: {exc}")
        return {"text": "", "source": "storage_not_found"}
    except urllib.error.URLError as exc:
        print(f"[extract-text] storage fetch failed for {name}: {exc}")
        return {"text": "", "source": "storage_fetch_failed"}
    except Exception as exc:
        print(f"[extract-text] read failed for {name}: {exc}")
        return {"text": "", "source": "read_error"}

    if len(raw) > MAX_EXTRACT_FILE_SIZE:
        return {"text": "", "source": "skipped_too_large"}

    try:
        text = _extract_text_from_bytes(
            raw, mime_type, page_from=page_from, page_to=page_to, name=name
        )
    except Exception as exc:
        print(f"[extract-text] parse failed for {name}: {exc}")
        return {"text": "", "source": "parse_error"}

    char_limit = max_chars if max_chars and max_chars > 0 else MAX_EXTRACTED_CHARS
    truncated = _truncate(text, char_limit)
    if truncated.endswith("…") and len(text) > char_limit:
        truncated = truncated[:-1] + " [truncated]"
    return {"text": truncated, "source": "ok"}


@app.post("/extract-text")
async def extract_text(body: ExtractRequest):
    """提取附件文本：data URL 走 base64；HTTP 路径走 fetch；按 mimeType 路由解析器。"""
    try:
        result = await asyncio.wait_for(
            asyncio.to_thread(
                extract_attachment_text,
                body.url,
                body.mimeType,
                body.name,
                body.size,
                body.page_from,
                body.page_to,
                body.max_chars,
            ),
            timeout=EXTRACT_TEXT_TIMEOUT_SECONDS,
        )
    except asyncio.TimeoutError:
        return JSONResponse(
            {"text": "", "source": "timeout", "name": body.name},
            status_code=200,
        )
    except Exception as exc:
        print(f"[extract-text] handler failed for {body.name}: {exc}")
        return JSONResponse(
            {"text": "", "source": "handler_error", "name": body.name},
            status_code=200,
        )

    return JSONResponse({**result, "name": body.name})
